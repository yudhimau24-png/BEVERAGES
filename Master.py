import io
import json
import os
import random
import re
import sqlite3
import hashlib
import time
import sys
import urllib.parse
import urllib.request
import requests
import pandas as pd
from PIL import Image
from pydantic import BaseModel, Field
import streamlit as st

from google import genai
from google.genai import types

# Selenium Imports untuk Robot Scraper
try:
    from selenium import webdriver
    from selenium.webdriver.common.by import By
    from selenium.webdriver.common.keys import Keys
    from selenium.webdriver.support.ui import WebDriverWait
    from selenium.webdriver.support import expected_conditions as EC
    from selenium.webdriver.support.ui import Select
    SELENIUM_AVAILABLE = True
except ImportError:
    SELENIUM_AVAILABLE = False

# Module openpyxl untuk Excel
from openpyxl import Workbook
from openpyxl.chart import BarChart, PieChart, Reference
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter

# Module python-pptx untuk Slide Presentasi
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Import dotenv secara aman
try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass


# ==============================================================================
# 0. DATABASE & LOCAL BPOM INGESTION SYSTEM (SQLITE)
# ==============================================================================
DB_FILE = "app_users.db"
FREE_USAGE_LIMIT = 3

def init_db():
    conn = sqlite3.connect(DB_FILE)
    c = conn.cursor()
    
    # Tabel Users & Authentication
    c.execute('''
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            full_name TEXT NOT NULL,
            email TEXT NOT NULL,
            phone TEXT,
            company TEXT,
            username TEXT UNIQUE NOT NULL,
            password_hash TEXT NOT NULL,
            role TEXT DEFAULT 'free',
            usage_count INTEGER DEFAULT 0,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    
    c.execute("PRAGMA table_info(users)")
    existing_columns = [column[1] for column in c.fetchall()]
    if "created_at" not in existing_columns:
        try:
            c.execute("ALTER TABLE users ADD COLUMN created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP")
        except Exception:
            pass
    
    c.execute("SELECT * FROM users WHERE username = 'admin'")
    if not c.fetchone():
        c.execute('''
            INSERT INTO users (full_name, email, phone, company, username, password_hash, role, usage_count)
            VALUES ('System Administrator', 'admin@beveragepro.com', '081234567890', 'BeveragePro HQ', 'admin', ?, 'admin', 0)
        ''', (hash_password('admin123'),))

    # TABEL DATASET LOKAL CEKBPOM RI
    c.execute('''
        CREATE TABLE IF NOT EXISTS bpom_records (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            tipe TEXT NOT NULL,
            nomor_registrasi TEXT UNIQUE NOT NULL,
            tanggal_terbit TEXT,
            nama_produk TEXT NOT NULL,
            merk TEXT NOT NULL,
            kemasan TEXT,
            pendaftar TEXT NOT NULL,
            lokasi TEXT
        )
    ''')
        
    conn.commit()
    conn.close()

def hash_password(password: str) -> str:
    return hashlib.sha256(password.encode()).hexdigest()

def register_user(full_name, email, phone, company, username, password):
    conn = sqlite3.connect(DB_FILE)
    c = conn.cursor()
    try:
        c.execute('''
            INSERT INTO users (full_name, email, phone, company, username, password_hash, role, usage_count)
            VALUES (?, ?, ?, ?, ?, ?, 'free', 0)
        ''', (full_name, email, phone, company, username, hash_password(password)))
        conn.commit()
        conn.close()
        return True, "Registrasi berhasil! Data Anda telah tersimpan di Database. Silakan Login."
    except sqlite3.IntegrityError:
        conn.close()
        return False, "Username sudah terdaftar! Pilih username lain."

def authenticate_user(username, password):
    conn = sqlite3.connect(DB_FILE)
    c = conn.cursor()
    c.execute("SELECT username, password_hash, role, usage_count, full_name FROM users WHERE username = ?", (username,))
    user = c.fetchone()
    conn.close()
    if user and user[1] == hash_password(password):
        return {
            "username": user[0],
            "role": user[2],
            "usage_count": user[3],
            "full_name": user[4]
        }
    return None

def get_user_data(username):
    conn = sqlite3.connect(DB_FILE)
    c = conn.cursor()
    c.execute("SELECT username, role, usage_count, full_name, email, phone, company, created_at FROM users WHERE username = ?", (username,))
    user = c.fetchone()
    conn.close()
    if user:
        return {
            "username": user[0],
            "role": user[1],
            "usage_count": user[2],
            "full_name": user[3],
            "email": user[4],
            "phone": user[5],
            "company": user[6],
            "created_at": user[7]
        }
    return None

def increment_user_usage(username):
    conn = sqlite3.connect(DB_FILE)
    c = conn.cursor()
    c.execute("UPDATE users SET usage_count = usage_count + 1 WHERE username = ?", (username,))
    conn.commit()
    conn.close()

def update_user_role(username, new_role):
    conn = sqlite3.connect(DB_FILE)
    c = conn.cursor()
    c.execute("UPDATE users SET role = ? WHERE username = ?", (new_role, username))
    conn.commit()
    conn.close()

def get_all_users_df():
    conn = sqlite3.connect(DB_FILE)
    df = pd.read_sql_query("SELECT id, full_name, email, phone, company, username, role, usage_count, created_at FROM users ORDER BY id DESC", conn)
    conn.close()
    return df

def delete_user_from_db(username):
    conn = sqlite3.connect(DB_FILE)
    c = conn.cursor()
    c.execute("DELETE FROM users WHERE username = ?", (username,))
    conn.commit()
    conn.close()


# ==============================================================================
# INTEGRASI SELENIUM AUTOMATION SCRAPER (MENAMPIKLAN UTUH 4 KOLOM)
# ==============================================================================
def run_selenium_bpom_robot(keyword: str):
    """Menyedot SELURUH data produk, merk, maupun Pendaftar/Sarana tanpa memotong kolom"""
    if not SELENIUM_AVAILABLE:
        return False, None, "Package Selenium belum terinstall di server!"
        
    try:
        options = webdriver.ChromeOptions()
        options.add_argument("--headless=new")
        options.add_argument("--no-sandbox")
        options.add_argument("--disable-dev-shm-usage")
        
        driver = webdriver.Chrome(options=options)
        driver.get("https://cekbpom.pom.go.id/")
        wait = WebDriverWait(driver, 15)

        all_collected_data = []

        daftar_kategori = [
            ("PANGAN OLAHAN", "//a[contains(., 'Pangan Olahan')] | //span[contains(., 'Pangan Olahan')] | //*[contains(text(), 'Pangan Olahan')]"),
            ("PRODUK", "//a[contains(., 'Produk')] | //span[contains(., 'Produk')]"),
            ("SARANA", "//a[contains(., 'Sarana')] | //span[contains(., 'Sarana')]"),
        ]

        js_extract = """
            var data = [];
            var rows = document.querySelectorAll('tbody tr');
            rows.forEach(function(tr) {
                var cols = Array.from(tr.querySelectorAll('td, th')).map(function(td) {
                    var text = td.innerText || "";
                    return text.split('\\n').join(' ').split('\\r').join(' ').trim();
                });
                if (cols.length > 1 && !cols[0].toLowerCase().includes('tidak ada data') && !cols[0].toLowerCase().includes('no data')) {
                    data.push(cols);
                }
            });
            return data;
        """

        for idx, (nama_kategori, xpath_menu) in enumerate(daftar_kategori):
            if idx > 0:
                try:
                    menu_btn = driver.find_element(By.XPATH, xpath_menu)
                    driver.execute_script("arguments[0].click();", menu_btn)
                    time.sleep(2.5)
                except Exception:
                    pass

            search_types = ["Pendaftar", "Merk", "Produk", "Sarana"]
            for s_type in search_types:
                try:
                    try:
                        filter_btn = wait.until(
                            EC.element_to_be_clickable((
                                By.XPATH,
                                "//button[contains(text(), 'Filter')] | //a[contains(text(), 'Filter')] | //*[contains(@class, 'filter')]"
                            ))
                        )
                        driver.execute_script("arguments[0].click();", filter_btn)
                        time.sleep(1)
                    except Exception:
                        pass

                    select_elements = driver.find_elements(By.TAG_NAME, "select")
                    for sel_elem in select_elements:
                        try:
                            select_obj = Select(sel_elem)
                            for opt in select_obj.options:
                                if s_type.lower() in opt.text.lower():
                                    select_obj.select_by_visible_text(opt.text)
                                    time.sleep(0.5)
                                    break
                        except Exception:
                            pass

                    input_box = wait.until(
                        EC.presence_of_element_located((
                            By.CSS_SELECTOR,
                            "input[placeholder*='Cari'], input[type='search'], input[type='text']"
                        ))
                    )
                    input_box.clear()
                    input_box.send_keys(keyword)
                    input_box.send_keys(Keys.ENTER)

                    time.sleep(3.5)

                    page = 1
                    while True:
                        rows = driver.execute_script(js_extract)
                        if rows:
                            for r in rows:
                                if r not in all_collected_data:
                                    all_collected_data.append(r)

                        next_btns = driver.find_elements(By.XPATH, "//a[contains(text(), 'Selanjutnya')] | //button[contains(text(), 'Selanjutnya')] | //li[contains(@class, 'next') and not(contains(@class, 'disabled'))]/a")
                        if not next_btns:
                            break

                        next_btn = next_btns[0]
                        parent_class = next_btn.find_element(By.XPATH, "..").get_attribute("class") or ""
                        btn_class = next_btn.get_attribute("class") or ""

                        if "disabled" in parent_class.lower() or "disabled" in btn_class.lower() or next_btn.get_attribute("disabled") is not None:
                            break

                        old_first = rows[0][0] if rows else ""
                        driver.execute_script("arguments[0].click();", next_btn)
                        time.sleep(2.5)

                        new_rows = driver.execute_script(js_extract)
                        new_first = new_rows[0][0] if new_rows else ""
                        if old_first == new_first:
                            break

                        page += 1
                        if page > 20:
                            break

                    if all_collected_data:
                        break
                except Exception:
                    continue

        driver.quit()

        if all_collected_data:
            df_hasil = pd.DataFrame(all_collected_data)
            num_cols = len(df_hasil.columns)
            
            # MEMASTIKAN TIDAK ADA KOLOM YANG DIPOTONG
            if num_cols == 4:
                df_hasil.columns = ['Tipe', 'Nomor Registrasi', 'Nama Produk (Merk)', 'Pendaftar / Sarana']
            elif num_cols == 3:
                df_hasil.columns = ['Nomor Registrasi', 'Nama Produk (Merk)', 'Pendaftar / Sarana']
            else:
                df_hasil.columns = [f"Kolom {i+1}" for i in range(num_cols)]

            return True, df_hasil, f"🎉 Berhasil mengambil TOTAL {len(df_hasil)} data produk & sarana/pendaftar untuk '{keyword}'!"

        return False, None, f"Data '{keyword}' tidak ditemukan di Pangan Olahan, Produk, maupun Sarana."

    except Exception as e:
        try: driver.quit()
        except: pass
        return False, None, f"Terjadi kendala pada Selenium Robot: {str(e)}"


# ==============================================================================
# SMART POSITIONAL & PATTERN BPOM INGESTION ENGINE
# ==============================================================================
def clean_search_keyword(kw: str) -> str:
    stopwords = [
        "semua", "daftar", "cari", "merk", "brand", "produk", "list", "sarana",
        "chateau", "château", "domaine", "domain", "tbk", "inc", "ltd"
    ]
    words = kw.strip().split()
    cleaned_words = [w for w in words if w.lower() not in stopwords]
    if cleaned_words:
        return " ".join(cleaned_words)
    return kw.strip()

def save_fetched_records_to_sqlite(records: list):
    records_tuple = []
    for r in records:
        tipe_code = r.get('tipe', 'PO').strip().upper()
        if tipe_code in ['MD', 'ML', 'PO', 'PANGAN OLAHAN']:
            tipe_code = 'PO'
            
        records_tuple.append((
            tipe_code,
            r.get('nomor_registrasi', '').strip(),
            r.get('tanggal_terbit', 'N/A').strip(),
            r.get('nama_produk', '').strip(),
            r.get('merk', 'N/A').strip().upper(),
            r.get('kemasan', 'N/A').strip(),
            r.get('pendaftar', 'N/A').strip().upper(),
            r.get('lokasi', 'Indonesia').strip()
        ))
    
    conn = sqlite3.connect(DB_FILE)
    c = conn.cursor()
    c.executemany('''
        INSERT OR REPLACE INTO bpom_records 
        (tipe, nomor_registrasi, tanggal_terbit, nama_produk, merk, kemasan, pendaftar, lokasi)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?)
    ''', records_tuple)
    conn.commit()
    conn.close()

def ingest_bpom_dataframe(df: pd.DataFrame):
    """Membaca & Mem-parsing Data ke Database SQLite secara Presisi"""
    if df.empty:
        return False, "Dataset kosong!"
        
    records = []
    df = df.astype(str)
    
    for _, row in df.iterrows():
        row_vals = [str(v).strip() for v in row.values if str(v).strip() and str(v).strip().lower() != 'nan']
        if not row_vals:
            continue
            
        combined_row_text = " ".join(row_vals).lower()
        if "nomor registrasi" in combined_row_text and "nama produk" in combined_row_text:
            continue

        raw_reg = ""
        raw_prod = ""
        raw_pend = ""
        raw_tipe = ""
        
        for val in row_vals:
            if re.search(r'(MD|ML|NA|TR|SD|SL)\s*\d+', val, re.IGNORECASE) or "terbit:" in val.lower():
                raw_reg = val
                break
                
        for val in row_vals:
            if "merk:" in val.lower() or "kemasan:" in val.lower() or (val != raw_reg and len(val) > 15 and not any(x in val.lower() for x in ['pt ', 'cv ', 'kota ', 'kab '])):
                raw_prod = val
                break

        for val in reversed(row_vals):
            if val != raw_reg and val != raw_prod and len(val) > 3:
                raw_pend = val
                break

        if not raw_reg and len(row_vals) >= 2: raw_reg = row_vals[1]
        if not raw_prod and len(row_vals) >= 3: raw_prod = row_vals[2]
        if not raw_pend and len(row_vals) >= 4: raw_pend = row_vals[3]
        if len(row_vals) >= 1 and str(row_vals[0]).strip().upper() in ['PO', 'KO', 'TR', 'MD', 'ML']:
            raw_tipe = str(row_vals[0]).strip().upper()

        no_reg = raw_reg
        tgl_terbit = "N/A"
        if "terbit:" in raw_reg.lower():
            parts = re.split(r'terbit:', raw_reg, flags=re.IGNORECASE)
            no_reg = parts[0].strip()
            if len(parts) > 1:
                tgl_terbit = parts[1].strip()

        if no_reg.startswith("ML") or no_reg.startswith("MD") or raw_tipe in ['PO', 'MD', 'ML', 'PANGAN OLAHAN']:
            tipe_clean = "PO"
        elif no_reg.startswith("NA") or no_reg.startswith("KO") or raw_tipe in ['KO', 'NA']:
            tipe_clean = "KO"
        elif no_reg.startswith("TR") or raw_tipe == 'TR':
            tipe_clean = "TR"
        else:
            tipe_clean = "PO"

        nama_produk = raw_prod
        merk_clean = "N/A"
        kemasan_clean = "N/A"

        match_merk = re.search(r'Merk:\s*([^Kemasan\n]+)', raw_prod, re.IGNORECASE)
        if match_merk:
            merk_clean = match_merk.group(1).strip().upper()

        match_kemasan = re.search(r'Kemasan:\s*(.+)', raw_prod, re.IGNORECASE)
        if match_kemasan:
            kemasan_clean = match_kemasan.group(1).strip()

        if "Merk:" in nama_produk:
            nama_produk = re.split(r'Merk:', nama_produk, flags=re.IGNORECASE)[0].strip()

        pendaftar_clean = raw_pend.upper()
        lokasi_clean = "Indonesia"
        if " KOTA " in pendaftar_clean or " KAB " in pendaftar_clean or " JAKARTA " in pendaftar_clean:
            match_lokasi = re.search(r'(KOTA|KAB|JAKARTA|PROVINSI).+', pendaftar_clean)
            if match_lokasi:
                lokasi_clean = match_lokasi.group(0).strip()
                pendaftar_clean = pendaftar_clean.replace(lokasi_clean, "").strip()

        if no_reg and no_reg.lower() not in ['nan', 'none', '']:
            records.append({
                "tipe": tipe_clean,
                "nomor_registrasi": no_reg,
                "tanggal_terbit": tgl_terbit,
                "nama_produk": nama_produk if nama_produk else "Produk Terdaftar",
                "merk": merk_clean if merk_clean != "N/A" else "BPOM REGISTERED",
                "kemasan": kemasan_clean,
                "pendaftar": pendaftar_clean if pendaftar_clean else "PT IMPORTIR / PRODUSEN",
                "lokasi": lokasi_clean
            })

    if not records:
        return False, "Tidak ada data valid yang bisa di-parse dari file tersebut!"

    save_fetched_records_to_sqlite(records)
    return True, f"🎉 Berhasil mengekstrak & menyimpan {len(records)} record BPOM ke database lokal!"

def search_bpom_sqlite_hybrid(raw_keyword: str, filter_tipe: str = "Semua", api_key: str = "", limit: int = 200):
    clean_kw_str = clean_search_keyword(raw_keyword)
    words = [w.strip().lower() for w in clean_kw_str.split() if len(w.strip()) > 0]
    
    conn = sqlite3.connect(DB_FILE)
    c = conn.cursor()
    
    if not words:
        conn.close()
        return [], "NOT_FOUND", raw_keyword, "Kata kunci pencarian kosong."

    combined_col = "(LOWER(COALESCE(nomor_registrasi,'')) || ' ' || LOWER(COALESCE(nama_produk,'')) || ' ' || LOWER(COALESCE(merk,'')) || ' ' || LOWER(COALESCE(pendaftar,'')) || ' ' || LOWER(COALESCE(kemasan,'')))"
    
    where_clauses = [f"{combined_col} LIKE ?" for _ in words]
    params = [f"%{w}%" for w in words]
    
    if filter_tipe != "Semua":
        if filter_tipe.upper() in ["PO", "MD", "ML"]:
            where_clauses.append("tipe IN ('PO', 'MD', 'ML')")
        else:
            where_clauses.append("tipe = ?")
            params.append(filter_tipe.upper())
        
    where_sql = " AND ".join(where_clauses)
    
    query = f'''
        SELECT tipe, nomor_registrasi, tanggal_terbit, nama_produk, merk, kemasan, pendaftar, lokasi
        FROM bpom_records
        WHERE {where_sql}
        ORDER BY id ASC
        LIMIT ?
    '''
    params.append(limit)
    
    c.execute(query, tuple(params))
    rows = c.fetchall()

    status_fallback = None
    if not rows:
        sig_words = [w for w in words if len(w) > 2]
        for sw in sig_words:
            fb_params = [f"%{sw}%"]
            fb_where = [f"{combined_col} LIKE ?"]
            if filter_tipe != "Semua":
                if filter_tipe.upper() in ["PO", "MD", "ML"]:
                    fb_where.append("tipe IN ('PO', 'MD', 'ML')")
                else:
                    fb_where.append("tipe = ?")
                    fb_params.append(filter_tipe.upper())
            
            fb_query = f'''
                SELECT tipe, nomor_registrasi, tanggal_terbit, nama_produk, merk, kemasan, pendaftar, lokasi
                FROM bpom_records
                WHERE {" AND ".join(fb_where)}
                ORDER BY id ASC
                LIMIT ?
            '''
            fb_params.append(limit)
            c.execute(fb_query, tuple(fb_params))
            rows = c.fetchall()
            if rows:
                status_fallback = f"ℹ️ Hasil disesuaikan untuk kata kunci utama **'{sw.upper()}'**:"
                break

    conn.close()

    if rows:
        results = []
        for r in rows:
            display_tipe = "PO" if r[0] in ["MD", "ML", "PO"] else r[0]
            results.append({
                "tipe": display_tipe, "nomor_registrasi": r[1], "tanggal_terbit": r[2],
                "nama_produk": r[3], "merk": r[4], "kemasan": r[5],
                "pendaftar": r[6], "lokasi": r[7]
            })
        return results, "SQLITE_LOCAL", clean_kw_str, status_fallback

    if not api_key:
        return [], "NO_API_KEY", raw_keyword, f"Data '{raw_keyword}' belum ada di database. Masukkan Gemini API Key di Sidebar untuk mencari otomatis!"

    try:
        client = genai.Client(api_key=api_key)
        prompt_bpom = f"""
        Ekstrak SELURUH VARIAN RESMI data CekBPOM RI (https://cekbpom.pom.go.id/) untuk merek/brand/sarana/pendaftar: '{clean_kw_str}'.

        ATURAN SANGAT KETAT:
        1. Sebutkan SEMUA varian resmi yang terdaftar untuk merek/pendaftar '{clean_kw_str}'.
        2. Tipe: Wajib tulis 'PO' untuk Pangan Olahan/Makanan/Minuman (MD/ML), 'KO' (Kosmetik), 'TR' (Obat Tradisional).
        3. Pendaftar: Tulis nama PT pendaftar/importir resmi di Indonesia (misal: PT JADDI INTERNASIONAL).
        4. Jika merek/pendaftar '{clean_kw_str}' TIDAK TERDAFTAR RESMI DI BPOM, KEMBALIKAN ARRAY KOSONG PERSIS: []

        Format Output JSON Array:
        [
            {{
                "tipe": "PO",
                "nomor_registrasi": "ML/MD xxxxxxxxxxxx",
                "tanggal_terbit": "YYYY-MM-DD",
                "nama_produk": "Deskripsi Produk Varian BPOM",
                "merk": "MERK PRODUK",
                "kemasan": "Bentuk Kemasan",
                "pendaftar": "{clean_kw_str.upper()}",
                "lokasi": "Kota/Kab, Provinsi, Indonesia"
            }}
        ]
        """
        res = client.models.generate_content(
            model="gemini-3.5-flash-lite",
            contents=prompt_bpom,
            config=types.GenerateContentConfig(temperature=0.0)
        )
        
        raw_text = res.text.strip()
        clean_json = re.sub(r'```json\s*|\s*```', '', raw_text).strip()
        match = re.search(r'\[.*\]', clean_json, re.DOTALL)
        
        if match:
            fetched_data = json.loads(match.group(0))
            if isinstance(fetched_data, list) and len(fetched_data) > 0:
                save_fetched_records_to_sqlite(fetched_data)
                
                filtered = [d for d in fetched_data if filter_tipe == "Semua" or d.get("tipe") in [filter_tipe, "PO", "MD", "ML"]]
                note_ai = None
                if not filtered and filter_tipe != "Semua":
                    filtered = fetched_data
                    note_ai = f"ℹ️ Catatan: '{raw_keyword.upper()}' terdeteksi sebagai produk kategori **{fetched_data[0].get('tipe')}**."
                    
                return filtered, "AI_AUTO_FETCHED", clean_kw_str, note_ai

        return [], "NOT_FOUND", raw_keyword, f"Data resmi CekBPOM RI tidak ditemukan untuk kata kunci '{raw_keyword}'."

    except Exception as e:
        err_str = str(e)
        if "429" in err_str or "RESOURCE_EXHAUSTED" in err_str:
            return [], "ERROR", raw_keyword, "⚠️ **KUOTA API GEMINI HABIS (Error 429)**: Silakan tunggu 1-2 menit atau ganti API Key di Sidebar."
        return [], "ERROR", raw_keyword, f"Terjadi kesalahan AI: {err_str}"

init_db()


# ==============================================================================
# 1. SOMMELIER MASTER REGISTRY & SANITIZER
# ==============================================================================
MASTER_VERIFIED_BRANDS = {
    "kawa": {
        "real_name": "Anggur Kawa-Kawa",
        "producer": "PT Balaraja Barat Indah (PT BBI)",
        "origin_type": "Lokal / Indonesia",
        "abv": "19.8%",
        "behind_story": "Diproduksi secara mandiri oleh PT Balaraja Barat Indah di Balaraja, Tangerang.",
    }
}

REAL_COMPETITOR_FALLBACKS = {
    "alco pops": [
        {
            "product_name": "Mix Max Exotic Flavors",
            "origin_brand": "PT Astidama Adimukti (Lokal / Indonesia)",
            "origin_type": "Lokal / Indonesia",
            "alcohol_by_volume": "4.8%",
            "base_ingredients": "Vodka netral, perisa buah, air berkarbonasi",
            "price_point": "IDR 25,000 - IDR 35,000",
            "key_difference": "Alco pop lokal dengan variasi rasa buah manis",
            "product_weakness": "Kandungan gula cukup tinggi",
            "behind_story": "Pionir RTD lokal buatan PT Astidama Adimukti."
        }
    ]
}

def HARD_LOCK_SANITIZER(data_json: dict) -> dict:
    beverages = data_json.get("selected_beverages", [])
    for bev in beverages:
        bev_name = bev.get("beverage_name", "")
        for key, truth in MASTER_VERIFIED_BRANDS.items():
            if key in bev_name.lower():
                bev["beverage_name"] = truth["real_name"]
                bev["origin_type"] = truth["origin_type"]
                bev["alcohol_by_volume"] = truth["abv"]
                bev["behind_story"] = truth["behind_story"]
                break
    return data_json

def clean_text(text: str) -> str:
    if not text:
        return ""
    cleaned = re.sub(r"[^\x00-\x7F]+", "", text)
    return re.sub(r"\s+", " ", cleaned).strip()


# ==========================================
# 2. PYDANTIC SCHEMAS
# ==========================================
class CompetitorItem(BaseModel):
    product_name: str = Field(description="Nama produk pembanding real")
    origin_brand: str = Field(description="Asal negara / brand produsen")
    origin_type: str = Field(description="Lokal / Indonesia atau Impor")
    alcohol_by_volume: str = Field(description="Kadar alkohol (% ABV)")
    base_ingredients: str = Field(description="Bahan baku")
    price_point: str = Field(description="Rentang harga")
    key_difference: str = Field(description="Perbedaan utama")
    product_weakness: str = Field(description="Kelemahan produk")
    behind_story: str = Field(description="Sejarah singkat")

class BusinessAnalysis(BaseModel):
    target_demographics: str = Field(description="Target konsumen")
    markup_margin_potential: str = Field(description="Potensi Margin")
    suitable_venue: str = Field(description="Venue terbaik")
    investment_value: str = Field(description="Potensi investasi")

class BeverageItem(BaseModel):
    beverage_name: str = Field(description="Nama minuman")
    beverage_category: str = Field(description="Kategori")
    origin_type: str = Field(description="Asal produk")
    alcohol_by_volume: str = Field(description="Kadar alkohol (% ABV)")
    unique_selling_point: str = Field(description="USP")
    behind_story: str = Field(description="Sejarah singkat")
    product_weakness: str = Field(description="Kelemahan")
    age_or_vintage: str = Field(description="Vintage")
    origin: str = Field(description="Negara & Region")
    base_ingredients: str = Field(description="Bahan baku")
    rating_scores: str = Field(description="Rating")
    estimated_price: str = Field(description="Harga")
    flavor_character: list[str] = Field(description="Notes")
    serving_recommendation: str = Field(description="Cara saji")
    business_intelligence: BusinessAnalysis = Field(description="Bisnis")
    equivalent_competitors: list[CompetitorItem] = Field(description="Kompetitor")
    expert_note: str = Field(description="Catatan pakar")
    image_url: str = Field(default="")

class SelectorResponse(BaseModel):
    executive_summary: str = Field(description="Ringkasan")
    selected_beverages: list[BeverageItem] = Field(description="Daftar minuman")


# ==========================================
# 3. EXCEL & PPT GENERATOR
# ==========================================
def create_excel_dashboard(data_json: dict) -> io.BytesIO:
    wb = Workbook()
    ws = wb.active
    ws.title = "Executive Dashboard"
    ws.views.sheetView[0].showGridLines = True

    title_font = Font(name="Calibri", size=16, bold=True, color="FFFFFF")
    title_fill = PatternFill(start_color="4C1D95", end_color="4C1D95", fill_type="solid")
    header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="6D28D9", end_color="6D28D9", fill_type="solid")
    sec_font = Font(name="Calibri", size=12, bold=True, color="FFFFFF")
    sec_fill = PatternFill(start_color="3B0764", end_color="3B0764", fill_type="solid")
    thin_border = Border(
        left=Side(style="thin", color="E9D5FF"), right=Side(style="thin", color="E9D5FF"),
        top=Side(style="thin", color="E9D5FF"), bottom=Side(style="thin", color="E9D5FF"),
    )

    ws.merge_cells("A1:T2")
    banner = ws["A1"]
    banner.value = "GLOBAL BEVERAGE INTELLIGENCE DASHBOARD"
    banner.font = title_font
    banner.fill = title_fill
    banner.alignment = Alignment(horizontal="center", vertical="center")

    ws.merge_cells("A4:T6")
    summary_cell = ws["A4"]
    exec_text = data_json.get("executive_summary", "")
    summary_cell.value = f"📌 EXECUTIVE SUMMARY:\n{exec_text}"
    summary_cell.font = Font(name="Calibri", size=10, italic=True)
    summary_cell.fill = PatternFill(start_color="F3E8FF", end_color="F3E8FF", fill_type="solid")
    summary_cell.alignment = Alignment(vertical="center", horizontal="left", wrap_text=True)

    ws_chart_data = wb.create_sheet(title="ChartData")
    ws_chart_data["A1"], ws_chart_data["B1"] = "Kategori", "Jumlah Produk"

    items = data_json.get("selected_beverages", [])
    cat_counts = {}
    for item in items:
        cat = item.get("beverage_category", "Lainnya")
        cat_counts[cat] = cat_counts.get(cat, 0) + 1

    r = 2
    for cat, count in cat_counts.items():
        ws_chart_data.cell(row=r, column=1, value=cat)
        ws_chart_data.cell(row=r, column=2, value=count)
        r += 1

    labels_ref = Reference(ws_chart_data, min_col=1, min_row=2, max_row=r - 1)
    data_ref = Reference(ws_chart_data, min_col=2, min_row=1, max_row=r - 1)

    pie = PieChart()
    pie.add_data(data_ref, titles_from_data=True)
    pie.set_categories(labels_ref)
    pie.title, pie.width, pie.height = "Proporsi Kategori", 14, 7.5
    ws.add_chart(pie, "A8")

    bar = BarChart()
    bar.type, bar.style = "col", 10
    bar.title, bar.y_axis.title = "Sebaran Produk", "Jumlah Botol"
    bar.add_data(data_ref, titles_from_data=True)
    bar.set_categories(labels_ref)
    bar.width, bar.height = 16, 7.5
    ws.add_chart(bar, "I8")

    start_row = 23
    ws.merge_cells(f"A{start_row-1}:T{start_row-1}")
    sec_title = ws[f"A{start_row-1}"]
    sec_title.value = "📊 DETAILED BEVERAGE ANALYSIS"
    sec_title.font, sec_title.fill = sec_font, sec_fill
    sec_title.alignment = Alignment(horizontal="left", vertical="center")

    headers = [
        "No", "Nama Minuman", "Kategori", "Klasifikasi", "ABV", "Vintage", "Asal",
        "Bahan Baku", "Behind Story", "Rating", "Harga", "USP", "Kelemahan",
        "Tasting Notes", "Cara Saji", "Pembanding", "Margin", "Target", "Venue", "Investasi"
    ]

    for col_idx, h in enumerate(headers, 1):
        cell = ws.cell(row=start_row, column=col_idx)
        cell.value, cell.font, cell.fill, cell.border = h, header_font, header_fill, thin_border
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)

    for idx, item in enumerate(items, 1):
        curr_row = start_row + idx
        biz = item.get("business_intelligence", {})
        comps = item.get("equivalent_competitors", [])

        comp_str = "\n".join([f"• {c.get('product_name')} ({c.get('origin_type')}) | ABV: {c.get('alcohol_by_volume')} | Story: {c.get('behind_story')}" for c in comps])

        row_values = [
            idx, item.get("beverage_name"), item.get("beverage_category"), item.get("origin_type"),
            item.get("alcohol_by_volume"), item.get("age_or_vintage"), item.get("origin"),
            item.get("base_ingredients"), item.get("behind_story"), item.get("rating_scores"),
            item.get("estimated_price"), item.get("unique_selling_point"), item.get("product_weakness"),
            ", ".join(item.get("flavor_character", [])), item.get("serving_recommendation"),
            comp_str, biz.get("markup_margin_potential"), biz.get("target_demographics"),
            biz.get("suitable_venue"), biz.get("investment_value"),
        ]

        for col_idx, val in enumerate(row_values, 1):
            cell = ws.cell(row=curr_row, column=col_idx)
            cell.value, cell.border = val, thin_border
            cell.alignment = Alignment(vertical="center", wrap_text=True)
            if idx % 2 == 0:
                cell.fill = PatternFill(start_color="FAF5FF", end_color="FAF5FF", fill_type="solid")

    for col in ws.columns:
        col_letter = get_column_letter(col[0].column)
        max_len = max([len(str(cell.value or "")) for cell in col if cell.row >= start_row] + [0])
        ws.column_dimensions[col_letter].width = min(max(max_len + 3, 12), 35)

    buffer = io.BytesIO()
    wb.save(buffer)
    buffer.seek(0)
    return buffer


def create_pptx_presentation(data_json: dict) -> io.BytesIO:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    theme = {"card_bg": RGBColor(15, 23, 42), "border": RGBColor(234, 179, 8), "primary": RGBColor(250, 204, 21), "text": RGBColor(241, 245, 249), "card_header": RGBColor(30, 41, 59)}

    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = theme["card_bg"]
    tf1 = bg1.text_frame
    p1 = tf1.paragraphs[0]
    p1.text, p1.font.size, p1.font.bold, p1.font.color.rgb = "GLOBAL BEVERAGE INTELLIGENCE", Pt(36), True, theme["primary"]

    for idx, item in enumerate(data_json.get("selected_beverages", []), 1):
        slide_a = prs.slides.add_slide(blank_layout)
        bg_a = slide_a.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg_a.fill.solid()
        bg_a.fill.fore_color.rgb = theme["card_bg"]

        card_info = slide_a.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.3))
        card_info.fill.solid()
        card_info.fill.fore_color.rgb = theme["card_header"]
        tf_info = card_info.text_frame
        tf_info.word_wrap = True

        p_spec = tf_info.paragraphs[0]
        p_spec.text = f"#{idx} {item.get('beverage_name')}\n\n• ABV: {item.get('alcohol_by_volume')} | Klasifikasi: {item.get('origin_type')}\n• Behind Story: {item.get('behind_story')}\n• USP: {item.get('unique_selling_point')}"
        p_spec.font.size, p_spec.font.color.rgb = Pt(14), theme["text"]

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


# ==========================================
# 4. STREAMLIT APP CONFIG & STYLING
# ==========================================
st.set_page_config(
    page_title="Global Beverage & Business AI Pro",
    page_icon="🍸",
    layout="wide",
    initial_sidebar_state="expanded"
)

def apply_custom_background(image_url: str):
    st.markdown(
        f"""
        <style>
        @import url('https://fonts.googleapis.com/css2?family=Plus+Jakarta+Sans:wght@400;500;600;700;800&display=swap');
        html, body, [class*="css"] {{ font-family: 'Plus Jakarta Sans', sans-serif !important; }}
        
        .stApp {{
            background: linear-gradient(rgba(11, 15, 25, 0.5), rgba(11, 15, 25, 0.7)), url("{image_url}");
            background-size: cover; background-position: center; background-attachment: fixed; color: #F8FAFC;
        }}
        .hero-container {{
            background: rgba(15, 23, 42, 0.80); backdrop-filter: blur(16px);
            border: 1px solid rgba(245, 158, 11, 0.4); border-radius: 20px; padding: 2rem; margin-bottom: 2rem;
        }}
        .hero-title {{
            font-size: 2.3rem; font-weight: 800;
            background: linear-gradient(135deg, #FDE047 0%, #F59E0B 100%);
            -webkit-background-clip: text; -webkit-text-fill-color: transparent;
        }}
        
        /* CSS UNTUK EFEK TOMBOL HIDUP (HOVER & ACTIVE) */
        div.stButton > button[kind="primary"] {{
            background: linear-gradient(135deg, #F59E0B 0%, #D97706 100%) !important;
            color: #0F172A !important; 
            font-weight: 800 !important; 
            border-radius: 12px !important;
            border: none !important;
            transition: all 0.3s ease-in-out !important;
            box-shadow: 0 4px 6px rgba(217, 119, 6, 0.3) !important;
        }}

        div.stButton > button[kind="primary"]:hover {{
            background: linear-gradient(135deg, #FCD34D 0%, #F59E0B 100%) !important;
            box-shadow: 0 6px 15px rgba(245, 158, 11, 0.6) !important;
            transform: translateY(-2px) !important; 
            color: #000000 !important;
        }}

        div.stButton > button[kind="primary"]:active {{
            background: linear-gradient(135deg, #D97706 0%, #B45309 100%) !important;
            box-shadow: 0 2px 4px rgba(217, 119, 6, 0.4) !important;
            transform: translateY(1px) !important; 
        }}

        /* EXACT CEKBPOM OFFICIAL TABLE UI STYLING */
        .bpom-card {{
            background-color: #FFFFFF !important;
            border: 1px solid #E2E8F0 !important;
            border-radius: 6px !important;
            padding: 1.1rem !important;
            margin-bottom: 0.8rem !important;
            color: #1E293B !important;
            box-shadow: 0 2px 8px rgba(0,0,0,0.08) !important;
            height: 100% !important;
        }}
        .bpom-tipe {{
            font-size: 1.1rem !important;
            font-weight: 800 !important;
            color: #1E293B !important;
            text-align: center !important;
            padding-top: 0.5rem !important;
        }}
        .bpom-reg {{
            font-size: 1.02rem !important;
            font-weight: 800 !important;
            color: #0F172A !important;
            margin-bottom: 0.2rem !important;
        }}
        .bpom-date {{
            font-size: 0.82rem !important;
            color: #64748B !important;
        }}
        .bpom-title {{
            font-size: 0.98rem !important;
            font-weight: 800 !important;
            color: #1E3A8A !important;
            margin-bottom: 0.4rem !important;
            line-height: 1.4 !important;
        }}
        .bpom-sub {{
            font-size: 0.88rem !important;
            color: #334155 !important;
            margin-bottom: 0.2rem !important;
        }}
        .bpom-pendaftar {{
            font-size: 0.95rem !important;
            font-weight: 800 !important;
            color: #0F172A !important;
            text-transform: uppercase !important;
        }}
        .bpom-lokasi {{
            font-size: 0.85rem !important;
            color: #475569 !important;
        }}
        </style>
        """,
        unsafe_allow_html=True
    )

THEME_IMAGE_URL = "https://images.unsplash.com/photo-1514933651103-005eec06c04b?q=80&w=1920&auto=format&fit=crop"
apply_custom_background(THEME_IMAGE_URL)


# ==============================================================================
# 5. ROUTING & UI CONTROLLER
# ==============================================================================
if "logged_in" not in st.session_state:
    st.session_state["logged_in"] = False
if "username" not in st.session_state:
    st.session_state["username"] = ""

if not st.session_state["logged_in"]:
    st.markdown('<div class="hero-container"><div class="hero-title">🍾 Global Beverage Business Intelligence Pro</div></div>', unsafe_allow_html=True)
    col_gate1, col_gate2, col_gate3 = st.columns([1, 2, 1])
    with col_gate2:
        tab_login, tab_signup = st.tabs(["🔑 LOGIN", "📝 SIGNUP"])
        with tab_login:
            login_user = st.text_input("Username:", key="login_username_input")
            login_pass = st.text_input("Password:", type="password", key="login_password_input")
            if st.button("🔓 LOG IN NOW", type="primary", use_container_width=True):
                user_auth = authenticate_user(login_user.strip(), login_pass.strip())
                if user_auth:
                    st.session_state["logged_in"] = True
                    st.session_state["username"] = user_auth["username"]
                    st.rerun()
                else:
                    st.error("❌ Username/Password salah!")
        with tab_signup:
            reg_name = st.text_input("Nama Lengkap:", key="reg_name")
            reg_email = st.text_input("Email:", key="reg_email")
            reg_user = st.text_input("Buat Username:", key="reg_user")
            reg_pass = st.text_input("Buat Password:", type="password", key="reg_pass")
            if st.button("📝 REGISTRASI AKUN", type="primary", use_container_width=True):
                if reg_name and reg_email and reg_user and reg_pass:
                    ok, msg = register_user(reg_name, reg_email, "", "", reg_user.strip(), reg_pass.strip())
                    st.success(msg) if ok else st.error(msg)

else:
    curr_user = get_user_data(st.session_state["username"])
    st.sidebar.markdown(f"### 👤 Profile: **{curr_user['full_name']}**")
    
    if curr_user['role'] == 'admin':
        st.sidebar.markdown("🛡️ **Akses Status:** `SYSTEM ADMINISTRATOR` 👑")
    elif curr_user['role'] == 'unlimited':
        st.sidebar.markdown("👑 **Akses Status:** `UNLIMITED TIER` 💎")
    else:
        rem_uses = max(0, FREE_USAGE_LIMIT - curr_user['usage_count'])
        st.sidebar.markdown(f"🐣 **Akses Status:** `FREE TIER` ({curr_user['usage_count']}/{FREE_USAGE_LIMIT} Terpakai)")
        st.sidebar.info(f"⏳ Sisa Akses Gratis: **{rem_uses} kali**")
        
        with st.sidebar.expander("💳 Upgrade ke Unlimited Tier"):
            st.write("Dapatkan akses analisis tanpa batas dengan upgrade akun Anda.")
            st.markdown("💰 **Biaya:** Rp 50.000 / Bulan")
            st.markdown("💳 **Transfer:** BCA `123-456-7890` a.n Sommelier Pro")
            if st.button("⚡ Simulasikan Pembayaran Berhasil", type="primary", key="sim_pay_btn"):
                update_user_role(curr_user['username'], 'unlimited')
                st.success("🎉 Pembayaran Dikonfirmasi! Akun Anda kini UNLIMITED!")
                st.rerun()

    if st.sidebar.button("🚪 Logout / Keluar", key="logout_btn"):
        st.session_state["logged_in"] = False
        st.session_state["username"] = ""
        st.rerun()

    st.sidebar.markdown("---")

    def get_auto_api_key() -> str:
        try:
            if "GEMINI_API_KEY" in st.secrets: return st.secrets["GEMINI_API_KEY"]
        except Exception: pass
        return os.environ.get("GEMINI_API_KEY", "")

    auto_api_key = get_auto_api_key()
    raw_api_key = st.sidebar.text_input("Gemini API Key:", type="password", value=auto_api_key)
    api_key = raw_api_key.strip()

    st.markdown('<div class="hero-container"><div class="hero-title">🍾 Global Beverage Business Intelligence Pro</div></div>', unsafe_allow_html=True)

    if curr_user['role'] == 'admin':
        tab_app, tab_scan, tab_admin = st.tabs(["BEVERAGES ANALYZE", "🔍 BPOM CHECK (HYBRID LOCAL & AI)", "🗄️ ADMIN PANEL"])
    else:
        tab_app, tab_scan = st.tabs(["BEVERAGES ANALYZE", "🔍 BPOM CHECK (HYBRID LOCAL & AI)"])
        tab_admin = None

    if 'selenium_msg' in st.session_state:
        ok_s, msg_s = st.session_state.pop('selenium_msg')
        if ok_s:
            st.balloons()
            st.toast("🎉 Robot Selenium Selesai Menyedot Data!", icon="✅")
            st.success(f"### 🎉 PROCESS COMPLETE!\n{msg_s}")
        else:
            st.error(msg_s)

    # TAB 1: MAIN SOMMELIER ANALYSIS APP
    with tab_app:
        selected_lang = st.radio(
            "Pilih bahasa laporan hasil analisis:",
            options=[
                "🌐 Dual Language (Indonesia & English)",
                "🇮🇩 Khusus Bahasa Indonesia",
                "🇬🇧 English Only"
            ],
            horizontal=True,
            key="app_language_selector"
        )
        st.markdown("<br>", unsafe_allow_html=True)

        with st.form("form_analisis"):
            col1, col2 = st.columns(2)
            with col1:
                beverage_type = st.selectbox(
                    "🧃 Beverages Category:",
                    [
                        "✨ Universal / Mix & Match Bebas (Campuran Terbaik)",
                        "🍹 Alco Pops / Ready To Drink ",
                        "🍷 Wine ",
                        "🥃 Whisky / Whiskey ",
                        "🍺 Bir / Craft Beer ",
                        "🍶 Sake ",
                        "🇰🇷 Soju ",
                        "🍸 Gin ",
                        "🧊 Vodka ",
                        "🌵 Tequila / Mezcal ",
                        "🏴‍☠️ Rum ",
                        "👑 Brandy "
                    ]
                )
                origin_scope = st.selectbox(
                    "🗺 Scope Of Origin:",
                    [
                        "🌐 Bebas Campur / Mix All (Lokal & Impor)",
                        "🇮🇩 Bangga Lokal / Asli Indonesia 🌴",
                        "🌍 Impor Premium / Luar Negeri 🚢",
                        "⚔️ Duel Head-to-Head (Lokal vs Impor) 🥊"
                    ]
                )
                brand_input = st.text_input("🏷️ Brand / Product (Optional):", placeholder="Contoh: ")
                item_count = st.selectbox("🔢 Jumlah Produk Yg Mau Dilihat:", options=list(range(1, 16)), index=2)
            with col2:
                sort_priority = st.selectbox(
                    "🎯 Main Priority:",
                    [
                        "🌈 Universal / Seimbang Semuanya",
                        "⭐ Rating Bintang Paling Tinggi",
                        "💰 Margin Profit Cuan Melimpah (HORECA)",
                        "⚖️ Best Value / Paling Worth It!"
                    ]
                )
                venue_type = st.selectbox(
                    "🏰 Vibes:",
                    [
                        "🎪 Universal / Masuk Semua Tipe Venue",
                        "🍸 Speakeasy Bar Rahasia & Cozy",
                        "🍷 Fine Dining Mewah & Romantis",
                        "🍻 Casual Pub Asyik Nongkrong",
                        "🪩 Nightclub Party Abis!"
                    ]
                )
                budget_category = st.selectbox(
                    "💸 Range Budget:",
                    [
                        "💳 Universal / Semua Range Price",
                        "🟢 Entry Level (< $30 - Ramah Kantong 🐣)",
                        "🟡 Premium Pour ($30 - $100 - Pas Buat Chill 🥂)",
                        "🟠 Top Shelf ($100 - $300 - Kelas Upper 👑)",
                        "🔴 Collector (> $300 - Level Sultan 💎)"
                    ]
                )
                origin_input = st.text_input("🌍 Specific Region (Optional):")

            btn_process = st.form_submit_button("🔍 Punch Down", type="primary", use_container_width=True)

        if btn_process:
            user_info = get_user_data(st.session_state["username"])
            
            if user_info['role'] not in ['unlimited', 'admin'] and user_info['usage_count'] >= FREE_USAGE_LIMIT:
                st.error("⚠️ EKSEKUSI DITOLAK: Quota Akses Gratis Anda Telah Habis (3/3)!")
                st.warning("Silakan upgrade ke UNLIMITED TIER di menu Sidebar untuk melanjutkan penggunaan tanpa batas.")
            elif not api_key:
                st.error("❌ API Key tidak terdeteksi atau kosong. Silakan atur di sidebar!")
            else:
                c_bev, c_scope, c_prio, c_ven, c_bud = clean_text(beverage_type), clean_text(origin_scope), clean_text(sort_priority), clean_text(venue_type), clean_text(budget_category)

                if "English" in selected_lang:
                    lang_instruction = "IMPORTANT: Generate ALL responses, descriptions, notes, and analysis strictly in ENGLISH."
                elif "Indonesia" in selected_lang:
                    lang_instruction = "PENTING: Hasilkan SEMUA respon, deskripsi, catatan, dan analisis sepenuhnya dalam BAHASA INDONESIA."
                else:
                    lang_instruction = (
                        "IMPORTANT / PENTING: Sajikan SEMUA teks, analisis, USP, kelemahan, behind story, dan catatan "
                        "dalam DUA BAHASA (Bahasa Indonesia dan English) secara berdampingan. "
                        "Format setiap kalimat/paragraf dengan pola: [Teks Bahasa Indonesia] / [English Text]."
                    )

                SYSTEM_INSTRUCTION = f"""
                Kamu adalah Global Beverages Director dan Master Sommelier Internasional. 
                KAMI AKAN MENGHAPUS SEMUA MERK FIKTIF. 
                Sebutkan sejarah pendirian, kelemahan produk, komposisi bahan baku, serta WAJIB memberikan 2 kompetitor setaranya secara riil. 
                Jika kamu tidak tahu kompetitor lokalnya, gunakan merk internasional yang sangat terkenal di pasaran sebagai pembandingnya.
                
                {lang_instruction}
                """

                prompt_query = f"Sortir {item_count} minuman: Kategori: {c_bev}, Scope Asal: {c_scope}, Brand/Acuan: {brand_input or 'Bebas'}, Prioritas: {c_prio}, Venue: {c_ven}, Budget: {c_bud}, Region: {origin_input or 'Bebas'}. Gunakan hanya merk riil."

                with st.spinner("📊 Analyze & Build Data(No Hallucination)..."):
                    try:
                        client = genai.Client(api_key=api_key)
                        try:
                            res = client.models.generate_content(
                                model="gemini-3.6-flash", contents=prompt_query,
                                config=types.GenerateContentConfig(system_instruction=SYSTEM_INSTRUCTION, temperature=0.1, response_mime_type="application/json", response_schema=SelectorResponse)
                            )
                        except Exception:
                            res = client.models.generate_content(
                                model="gemini-3.5-flash-lite", contents=prompt_query,
                                config=types.GenerateContentConfig(system_instruction=SYSTEM_INSTRUCTION, temperature=0.1, response_mime_type="application/json", response_schema=SelectorResponse)
                            )

                        raw_data = json.loads(res.text)
                        data = HARD_LOCK_SANITIZER(raw_data)

                        if user_info['role'] not in ['unlimited', 'admin']:
                            increment_user_usage(st.session_state["username"])

                        st.success(f"✅ Analisis & Benchmarking {len(data.get('selected_beverages', []))} Minuman Selesai!")

                        excel_buf, pptx_buf = create_excel_dashboard(data), create_pptx_presentation(data)

                        st.subheader("📥 Export Reports")
                        col_b1, col_b2 = st.columns(2)
                        col_b1.download_button("📊 Download Excel (.xlsx)", data=excel_buf, file_name="Executive_Beverage_Dashboard.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", type="primary", use_container_width=True)
                        col_b2.download_button("🖥️ Download PowerPoint (.pptx)", data=pptx_buf, file_name="Executive_Beverage_Presentation.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", type="primary", use_container_width=True)

                        st.markdown("---")
                        st.subheader("📋 Executive Summary")
                        st.info(data["executive_summary"])
                        st.markdown("---")

                        for i, bev in enumerate(data["selected_beverages"], 1):
                            biz = bev["business_intelligence"]
                            origin_tag = "🇮🇩 LOKAL" if "lokal" in bev.get("origin_type", "").lower() or "indonesia" in bev.get("origin_type", "").lower() else "🌍 IMPOR"

                            st.markdown(f"### #{i} {bev['beverage_name']} ({bev['age_or_vintage']}) &nbsp; `{origin_tag}` &nbsp; 🧪 `{bev['alcohol_by_volume']}`")
                            c1, c2, c3 = st.columns([1.5, 2, 2])
                            
                            c1.markdown(f"**🥃 Kategori / Category:** {bev['beverage_category']}\n\n**🏷️ Klasifikasi / Type:** `{bev['origin_type']}`\n\n**🧪 Kadar Alkohol / ABV:** `{bev['alcohol_by_volume']}`\n\n**🌾 Bahan Baku / Ingredients:** {bev['base_ingredients']}\n\n**📍 Asal / Origin:** {bev['origin']}\n\n**⭐ Rating:** `{bev['rating_scores']}`\n\n**💵 Harga / Price:** `{bev['estimated_price']}`")
                            c2.markdown(f"**📖 Behind Story:** {bev['behind_story']}\n\n**🔥 USP:** {bev['unique_selling_point']}\n\n**⚠️ Kelemahan / Weakness:** {bev['product_weakness']}\n\n**👃 Tasting Notes:** {', '.join(bev['flavor_character'])}\n\n**🧊 Saji / Serving:** {bev['serving_recommendation']}")
                            c3.markdown(f"**📈 Margin Potential:** {biz['markup_margin_potential']}\n\n**👥 Target Demographics:** {biz['target_demographics']}\n\n**🏢 Suitable Venue:** {biz['suitable_venue']}")

                            st.markdown("#### ⚔️ Produk Pembanding & Benchmarking Setara / Equivalent Competitors")
                            comps = bev.get("equivalent_competitors", [])
                            if comps:
                                comp_cols = st.columns(len(comps))
                                for idx_c, comp in enumerate(comps):
                                    c_tag = "🇮🇩 LOKAL" if "lokal" in comp.get("origin_type", "").lower() or "indonesia" in comp.get("origin_type", "").lower() else "🌍 IMPOR"
                                    comp_cols[idx_c].info(f"**{comp['product_name']}** `{c_tag}`\n\n📖 **Behind Story:** {comp['behind_story']}\n\n🧪 **ABV:** `{comp['alcohol_by_volume']}`\n\n📍 **Brand:** {comp['origin_brand']}\n\n🌾 **Ingredients:** {comp['base_ingredients']}\n\n💵 **Price:** `{comp['price_point']}`\n\n💡 **Difference:** {comp['key_difference']}\n\n⚠️ **Weakness:** {comp['product_weakness']}")
                            st.markdown("---")

                    except Exception as e:
                        st.error(f"Terjadi kesalahan / Error: {e}")

    # TAB 2: HYBRID AUTO-FETCH & AUTO-SAVE CEKBPOM SYSTEM
    with tab_scan:
        st.markdown("### 🔍 Search Database CekBPOM RI (Hybrid Auto-Fetch)")
        st.caption("Verifikasi Resmi Database CekBPOM: Cek SQLite -> AI Auto-Fetch -> Auto-Save ke Database Lokal")

        col_input1, col_input2 = st.columns([3, 1])
        with col_input1:
            target_brand_name = st.text_input("Cari Produk / Merk / Pendaftar:", placeholder="Contoh: Jaddi, PT Jaddi Internasional, Jameson, Terrazas...", key="manual_brand_input_bpom")
        with col_input2:
            filter_tipe_bpom = st.selectbox("Filter Tipe BPOM:", ["Semua", "PO", "KO", "TR"], index=0)

        scan_limit = st.number_input("Batas Tampilan Record:", min_value=1, max_value=200, value=200, key="scan_limit_key")
        
        col_act1, col_act2 = st.columns([2, 1])
        with col_act1:
            btn_do_scan = st.button("🔎 Filter & Cari Data BPOM", type="primary", key="btn_do_scan_bpom")
        with col_act2:
            btn_do_selenium = st.button("🤖 Run Selenium Bot(Direct Web)", key="btn_selenium_run")

        # PREVIEW DATA HASIL SCRAPING DENGAN DATAFRAME PRESISI 4 KOLOM
        if 'preview_selenium_df' in st.session_state:
            preview_df = st.session_state['preview_selenium_df']
            kw_prev = st.session_state.get('preview_selenium_kw', '')
            msg_txt = st.session_state.get('selenium_msg_text', '')

            st.toast(f"🎉 Robot Selenium sukses menarik {len(preview_df)} data BPOM!", icon="🤖")

            st.markdown("---")
            col_head1, col_head2 = st.columns([3, 1])
            col_head1.subheader(f"🔍 Preview Data Hasil Robot Selenium: '{kw_prev.upper()}' ({len(preview_df)} Item)")
            col_head2.metric("Total Terambil", f"{len(preview_df)} Data")
            
            if msg_txt:
                st.success(f"✨ {msg_txt}")

            st.dataframe(preview_df, use_container_width=True)

            col_p1, col_p2 = st.columns([2, 1])
            with col_p1:
                if st.button("💾 Simpan Data Ini ke Database SQLite", type="primary", key="btn_save_preview_data"):
                    ok_ing, msg_ing = ingest_bpom_dataframe(preview_df)
                    if ok_ing:
                        st.balloons()
                        st.toast("✅ Data berhasil disimpan ke SQLite!", icon="💾")
                        st.success(f"✅ {msg_ing}")
                        del st.session_state['preview_selenium_df']
                        if 'preview_selenium_kw' in st.session_state: del st.session_state['preview_selenium_kw']
                        if 'selenium_msg_text' in st.session_state: del st.session_state['selenium_msg_text']
                        st.rerun()
                    else:
                        st.error(msg_ing)
            with col_p2:
                if st.button("❌ Batal & Hapus Preview", key="btn_cancel_preview_data"):
                    del st.session_state['preview_selenium_df']
                    if 'preview_selenium_kw' in st.session_state: del st.session_state['preview_selenium_kw']
                    if 'selenium_msg_text' in st.session_state: del st.session_state['selenium_msg_text']
                    st.rerun()
            st.markdown("---")

        if btn_do_scan:
            if not target_brand_name.strip():
                st.warning("⚠️ Silakan ketik nama produk, merk, atau pendaftar terlebih dahulu!")
            else:
                st.markdown("---")
                with st.spinner(f"Memeriksa Database BPOM untuk '{target_brand_name}'..."):
                    bpom_records, status_code, clean_kw, note_fallback = search_bpom_sqlite_hybrid(target_brand_name.strip(), filter_tipe_bpom, api_key, scan_limit)

                    if status_code == "ERROR":
                        st.error(note_fallback)
                    elif status_code == "NO_API_KEY":
                        st.warning(note_fallback)
                    elif status_code == "NOT_FOUND" or not bpom_records:
                        st.warning(f"❌ Data resmi untuk '{clean_kw.upper()}' tidak ditemukan di CekBPOM RI.")
                        encoded_kw = urllib.parse.quote(clean_kw)
                        st.markdown(f"👉 [**Klik di sini untuk Cek Langsung ke Situs Resmi cekbpom.pom.go.id**](https://cekbpom.pom.go.id/search/produk/1/{encoded_kw})")
                    else:
                        if note_fallback:
                            st.info(note_fallback)
                        elif status_code == "SELENIUM_ONLINE_FETCHED":
                            st.success(f"🌐 Data resmi '{clean_kw.upper()}' berhasil disedot dari situs CekBPOM Online & **otomatis tersimpan ke SQLite**!")
                        elif status_code == "AI_AUTO_FETCHED":
                            st.success(f"✨ Seluruh varian '{clean_kw.upper()}' baru saja diambil via AI & **otomatis disimpan (*auto-ingest*) ke SQLite lokal**!")
                        elif status_code == "SQLITE_LOCAL":
                            st.info(f"⚡ Menampilkan seluruh varian '{clean_kw.upper()}' ({len(bpom_records)} Record) langsung dari Database Bulk Ingestion / SQLite Lokal (100% Akurat).")

                        # FORMAT DATAFRAME PRESISI DENGAN 4 KOLOM LENGKAP TERMASUK PENDAFTAR / SARANA
                        table_data = []
                        for rec in bpom_records:
                            col_reg = rec.get('nomor_registrasi', '')
                            tgl = rec.get('tanggal_terbit', '')
                            if tgl and tgl not in ['N/A', 'nan', '']:
                                if 'terbit:' not in col_reg.lower():
                                    col_reg += f" Terbit: {tgl}"

                            col_prod = rec.get('nama_produk', '')
                            merk = rec.get('merk', '')
                            kemasan = rec.get('kemasan', '')

                            if merk and merk not in ['N/A', 'NAN', 'NONE', ''] and 'merk:' not in col_prod.lower():
                                col_prod += f" Merk: {merk}"
                            if kemasan and kemasan not in ['N/A', 'nan', 'none', ''] and 'kemasan:' not in col_prod.lower():
                                col_prod += f" Kemasan: {kemasan}"

                            col_pend = rec.get('pendaftar', '')
                            lokasi = rec.get('lokasi', '')
                            if lokasi and lokasi not in ['N/A', 'Indonesia', 'nan', 'none', '']:
                                if lokasi.lower() not in col_pend.lower():
                                    col_pend += f" {lokasi}"

                            table_data.append({
                                "Tipe": rec.get('tipe', 'PO'),
                                "Nomor Registrasi": col_reg,
                                "Nama Produk (Merk)": col_prod,
                                "Pendaftar / Sarana": col_pend
                            })

                        df_result_display = pd.DataFrame(table_data)
                        st.dataframe(df_result_display, use_container_width=True)

        if btn_do_selenium:
            if not target_brand_name.strip():
                st.warning("⚠️ Masukkan kata kunci di kolom atas terlebih dahulu!")
            else:
                with st.spinner(f"🤖 Menjalankan Robot Selenium Scraper untuk '{target_brand_name}' langsung di web CekBPOM..."):
                    ok_sel, df_sel, msg_sel = run_selenium_bpom_robot(target_brand_name.strip())
                    if ok_sel and df_sel is not None:
                        st.session_state['preview_selenium_df'] = df_sel
                        st.session_state['preview_selenium_kw'] = target_brand_name.strip()
                        st.session_state['selenium_msg_text'] = msg_sel
                        st.rerun()
                    else:
                        st.error(msg_sel)

        st.markdown("<br><br>", unsafe_allow_html=True)
        with st.expander("📥 Bulk Data Ingestion: Upload Dataset CekBPOM Baru (CSV/Excel)"):
            st.markdown("### 🤖 Tool Bonus: Robot Scraper BPOM JS")
            st.write("Ingin mengambil ratusan data sekaligus dari situs BPOM secara otomatis via Browser?")
            st.info("1. Copy script di bawah ini.\n2. Buka situs [cekbpom.pom.go.id](https://cekbpom.pom.go.id) di tab baru.\n3. Tekan **F12**, pilih tab **Console**, lalu *Paste* dan tekan **Enter**!\n4. Robot akan bekerja mengunduh CSV untukmu, lalu upload hasilnya ke kolom di bawah.")
            
            st.code("""
            (async function robotBPOMAutoSearch() {
                let keyword = prompt("Masukkan kata kunci pencarian (misal: ML, 16.0, GULA, atau nama PT):");
                if (!keyword) return alert("Pencarian dibatalkan.");
                console.log(`🤖 Robot BPOM: Memulai pencarian untuk '${keyword}'...`);

                let inputCari = document.querySelector('input[placeholder*="Cari"], input[type="search"], input[type="text"]');
                let btnFilter = Array.from(document.querySelectorAll('button, input[type="submit"], a')).find(el => el.innerText.trim() === 'Filter');

                if (!inputCari || !btnFilter) return alert("❌ Gagal menemukan form pencarian!");

                inputCari.value = keyword;
                inputCari.dispatchEvent(new Event('input', { bubbles: true }));
                inputCari.dispatchEvent(new Event('change', { bubbles: true }));
                btnFilter.click();
                await new Promise(r => setTimeout(r, 3500));

                let allRows = [];
                let page = 1;

                while (true) {
                    let rows = Array.from(document.querySelectorAll('tbody tr'));
                    let pageData = rows.map(tr => Array.from(tr.querySelectorAll('td, th')).map(td => '"' + td.innerText.replace(/[\\r\\n]+/g, ' ').replace(/"/g, '""').trim() + '"').join(",")).filter(line => line.length > 5);

                    allRows.push(...pageData);

                    let nextBtn = Array.from(document.querySelectorAll('a, button, span, div')).find(el => el.innerText.trim() === 'Selanjutnya');
                    if (!nextBtn || nextBtn.classList.contains('disabled') || nextBtn.getAttribute('disabled') !== null) break;

                    let oldContent = document.querySelector('tbody')?.innerText;
                    nextBtn.click();
                    await new Promise(r => setTimeout(r, 2500));
                    
                    if (oldContent && oldContent === document.querySelector('tbody')?.innerText) break;
                    page++;
                }

                if (allRows.length === 0) return alert(`❌ Tidak ada data untuk '${keyword}'.`);

                let uniqueRows = Array.from(new Set(allRows));
                let csvContent = "\\uFEFF" + uniqueRows.join("\\n");
                let a = document.createElement("a");
                a.href = URL.createObjectURL(new Blob([csvContent], { type: 'text/csv;charset=utf-8;' }));
                a.download = `Data_BPOM_${keyword.replace(/[^a-zA-Z0-9]/g, '_')}.csv`;
                document.body.appendChild(a); a.click(); document.body.removeChild(a);
                alert(`🎉 BERHASIL! ${uniqueRows.length} data untuk '${keyword}' telah diunduh!`);
            })();
            """, language="javascript")
            
            st.markdown("---")
            st.write("Upload file dataset BPOM dalam format CSV atau Excel (.xlsx) hasil download Robot Scraper di atas ke dalam aplikasi:")
            
            uploaded_file = st.file_uploader("Pilih file CSV/Excel BPOM:", type=["csv", "xlsx"])
            if uploaded_file is not None:
                try:
                    if uploaded_file.name.endswith(".csv"):
                        df_upload = pd.read_csv(uploaded_file, header=None)
                    else:
                        df_upload = pd.read_excel(uploaded_file, header=None)
                        
                    st.dataframe(df_upload.head(3), use_container_width=True)
                    if st.button("⚡ Mulai Import / Ingest ke SQLite", type="primary"):
                        ok, msg = ingest_bpom_dataframe(df_upload)
                        if ok:
                            st.success(msg)
                            st.rerun()
                        else:
                            st.error(msg)
                except Exception as e_up:
                    st.error(f"Gagal membaca file: {e_up}")

    # TAB 3: ADMIN PANEL
    if curr_user['role'] == 'admin' and tab_admin is not None:
        with tab_admin:
            st.markdown("### 🗄️ Database User Management (Admin Dashboard)")
            users_df = get_all_users_df()
            
            # Summary Metrics Cards
            m1, m2, m3, m4 = st.columns(4)
            m1.metric("Total User Signup", len(users_df))
            m2.metric("Free Tier User", len(users_df[users_df['role'] == 'free']))
            m3.metric("Unlimited Tier User", len(users_df[users_df['role'] == 'unlimited']))
            m4.metric("Total Eksekusi AI", int(users_df['usage_count'].sum()))
            
            st.markdown("---")
            st.markdown("#### 📝 Edit & Sort User Role Langsung di Tabel:")
            st.caption("💡 Klik pada header kolom mana saja untuk melakukan **Sorting** (Urutkan Data), dan klik dropdown di kolom **role** untuk mengubah status akses user!")

            edited_users_df = st.data_editor(
                users_df,
                column_config={
                    "role": st.column_config.SelectboxColumn(
                        "role (Status Akses)",
                        help="Ubah role akses user",
                        width="medium",
                        options=["free", "unlimited", "admin"],
                        required=True,
                    )
                },
                disabled=["id", "full_name", "email", "phone", "company", "username", "usage_count", "created_at"],
                hide_index=True,
                use_container_width=True,
                key="admin_interactive_table"
            )

            col_btn1, col_btn2 = st.columns([2, 1])
            with col_btn1:
                if st.button("🔄 Simpan Perubahan Role", type="primary", key="save_role_btn_editor"):
                    changes_count = 0
                    for idx, row in edited_users_df.iterrows():
                        username_item = row['username']
                        new_r = row['role']
                        
                        old_r = users_df.loc[users_df['username'] == username_item, 'role'].values[0]
                        if new_r != old_r:
                            update_user_role(username_item, new_r)
                            changes_count += 1
                            
                    if changes_count > 0:
                        st.success(f"🎉 Berhasil memperbarui {changes_count} status role user ke database!")
                        st.rerun()
                    else:
                        st.info("ℹ️ Tidak ada perubahan role yang terdeteksi. Silakan ubah dropdown 'role' pada tabel di atas terlebih dahulu.")

            with col_btn2:
                csv_data = users_df.to_csv(index=False).encode('utf-8')
                st.download_button(
                    "📥 Export Database CSV",
                    data=csv_data,
                    file_name="User_Database_Signups.csv",
                    mime="text/csv",
                    use_container_width=True
                )
            
            st.markdown("---")
            st.markdown("#### ❌ Hapus Akun User")
            non_admin_users = users_df[users_df['username'] != 'admin']['username'].tolist()
            if non_admin_users:
                col_del1, col_del2 = st.columns([3, 1])
                with col_del1:
                    del_user = st.selectbox("Pilih Username yang Akan Dihapus:", options=non_admin_users, key="del_u_select")
                with col_del2:
                    st.markdown("<br>", unsafe_allow_html=True)
                    if st.button("❌ Hapus User", use_container_width=True):
                        delete_user_from_db(del_user)
                        st.success(f"User **{del_user}** telah dihapus dari database!")
                        st.rerun()
            else:
                st.info("Tidak ada pengguna lain selain admin utama.")
